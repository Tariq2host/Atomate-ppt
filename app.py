import streamlit as st
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import json
import io
import os
import win32com.client
import tempfile
import time
import pythoncom


FOLDER_RACINE = os.path.dirname(os.path.abspath(__file__))

def export_slide_with_shape_labels(ppt_path, output_folder, width=1280, height=720):
    import time
    from pptx import Presentation
    from PIL import Image, ImageDraw, ImageFont
    import win32com.client
    import os

    os.makedirs(output_folder, exist_ok=True)

    start_total = time.time()

    # Initialize COM for this thread
    pythoncom.CoInitialize()
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = True
        powerpoint.WindowState = 2  # minimized
        presentation = powerpoint.Presentations.Open(ppt_path, WithWindow=True)

        prs = Presentation(ppt_path)  # for shape info
        png_paths = []

        for i, slide in enumerate(presentation.Slides, start=1):
            start_slide = time.time()
            img_path = os.path.join(output_folder, f"slide_{i}.png")
            slide.Export(img_path, "PNG", width, height)
            png_paths.append(img_path)

            # Overlay shape names
            img = Image.open(img_path)
            draw = ImageDraw.Draw(img)
            try:
                font = ImageFont.truetype("arial.ttf", 14)
            except:
                font = ImageFont.load_default()

            slide_width_px, slide_height_px = width, height
            slide_width_emu, slide_height_emu = prs.slide_width, prs.slide_height
            scale_x = slide_width_px / slide_width_emu
            scale_y = slide_height_px / slide_height_emu

            for shape in prs.slides[i-1].shapes:
                if not shape.name:
                    continue
                if hasattr(shape, "left") and hasattr(shape, "top") and hasattr(shape, "width") and hasattr(shape, "height"):
                    x = int(shape.left * scale_x)
                    y = int(shape.top * scale_y)
                    w = int(shape.width * scale_x)
                    h = int(shape.height * scale_y)
                    draw.rectangle([x, y, x + w, y + h], outline="red", width=2)
                    draw.text((x + 3, y + 3), shape.name, fill="red", font=font)

            img.save(img_path)
            end_slide = time.time()
            print(f"Slide {i} processed in {end_slide - start_slide:.2f}s")

        presentation.Close()
        powerpoint.Quit()
        end_total = time.time()
        print(f"All slides processed in {end_total - start_total:.2f}s")
        return png_paths
    finally:
        pythoncom.CoUninitialize()



# --- Streamlit UI ---
st.set_page_config(page_title="PPT Mapping Assistant", layout="wide")
st.title("PowerPoint–Excel Mapping Assistant")
st.markdown("Build JSON mappings with slide previews showing shape names.")

st.header("📂 Upload your PowerPoint template")
ppt_file = st.file_uploader("Upload PowerPoint file (.pptx)", type=["pptx"])

if ppt_file:
    # Save uploaded PPTX temporarily
    tmp_ppt_path = os.path.join(tempfile.gettempdir(), "tmp_uploaded_ppt.pptx")
    with open(tmp_ppt_path, "wb") as f:
        f.write(ppt_file.getbuffer())
    ppt_path = tmp_ppt_path

    prs = Presentation(ppt_path)

    # Export slides with shape labels to fixed folder temp/img
    export_folder = os.path.join(FOLDER_RACINE, "temp", "img")
    print(export_folder)
    os.makedirs(export_folder, exist_ok=True)
    slide_images = export_slide_with_shape_labels(ppt_path, export_folder)

    # Build slide labels: number + title
    slide_labels = []
    for i, slide in enumerate(prs.slides, start=1):
        title_shape = next((s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()), None)
        title_text = title_shape.text.strip() if title_shape else "[No Title]"
        slide_labels.append(f"Slide {i} — {title_text[:50]}")

    # --- Slide selection ---
    selected_label = st.selectbox("Select a slide to map:", slide_labels)
    slide_number = int(selected_label.split("—")[0].replace("Slide", "").strip())
    slide = prs.slides[slide_number - 1]

    # --- Show slide preview ---
    st.markdown("### 🖼️ Slide Preview")
    if slide_images and len(slide_images) >= slide_number:
        st.image(slide_images[slide_number - 1], caption=f"Preview of Slide {slide_number}")
    else:
        st.warning("⚠ Slide preview not available.")

    st.markdown("---")
    st.subheader(f"🧩 Shapes found on Slide {slide_number}")
    st.write("Enter the corresponding Excel cell or range for each shape below:")

    # --- Shape mapping UI ---
    for idx, shape in enumerate(slide.shapes):
        if not shape.name:
            continue
        shape_type = "Text" if shape.has_text_frame else "Graph/Table"
        key = f"{slide_number}_{shape.name}_{idx}"
        col1, col2 = st.columns([2, 1])
        with col1:
            st.text_input(f"{shape.name} ({shape_type})", key=key, placeholder="e.g., A1 or Graphique1")
        with col2:
            st.markdown(f"*Type:* `{shape_type}`")

    # --- Generate mapping ---
    if st.button("💾 Generate Mapping"):
        text_mapping, graph_mapping = {}, {}
        for idx, shape in enumerate(slide.shapes):
            if not shape.name:
                continue
            key = f"{slide_number}_{shape.name}_{idx}"
            user_input = st.session_state.get(key)
            if not user_input:
                continue
            if shape.has_text_frame:
                text_mapping[shape.name] = {"cell": user_input}
            else:
                graph_mapping[shape.name] = user_input

        result = {"slide_index": slide_number, "text_mapping": text_mapping, "graph_mapping": graph_mapping}
        st.success(f"✅ Mapping generated for Slide {slide_number}!")
        st.json(result)

        json_str = json.dumps(result, indent=4)
        st.download_button(
            "📥 Download mapping as JSON",
            data=json_str.encode("utf-8"),
            file_name=f"ppt_excel_mapping_slide_{slide_number}.json",
            mime="application/json"
        )

else:
    st.info("👆 Please upload a PowerPoint file to start.")
